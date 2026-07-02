from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import re
from typing import List, Dict
from dataclasses import dataclass


@dataclass
class Finding:
    severity: str
    rule_id: str
    wcag_criterion: str
    message: str
    slide_number: int
    location: str
    fix_hint: str
    confidence: float = 1.0


class AccessibilityChecker:
    def __init__(self):
        self.rules = self._load_rules()

    def _load_rules(self):
        return [
            self._check_alt_text,
            self._check_slide_titles,
            self._check_empty_textboxes,
            self._check_duplicate_titles,           # New
            self._check_reading_order,              # New
            self._check_color_contrast,             # New
            self._check_very_small_text,            # New
            self._check_table_headers,              # New
            self._check_merged_table_cells,         # New
            self._check_image_only_slide,           # New
            self._check_chart_description,          # New
            self._check_text_embedded_in_image,     # New
            self._check_heading_structure,          # New
            self._check_text_overlapping,           # New
        ]

    def check_pptx(self, pptx_path: Path) -> List[Finding]:
        prs = Presentation(pptx_path)
        findings = []

        # Slide-level rules
        for slide_num, slide in enumerate(prs.slides, 1):
            for rule in self.rules:
                if rule.__name__ in ["_check_duplicate_titles"]:  # Skip if it's presentation-only
                    continue
                findings.extend(rule(slide, slide_num))

        # Presentation-level rules
        findings.extend(self._check_duplicate_titles(prs))

        return sorted(findings, key=lambda f: self._severity_order(f.severity), reverse=True)

    def _severity_order(self, sev: str) -> int:
        order = {"critical": 4, "serious": 3, "moderate": 2, "minor": 1}
        return order.get(sev, 0)

    # ========================= EXISTING RULES =========================

    def _check_alt_text(self, slide, slide_num: int) -> List[Finding]:
        findings = []
        for shape_idx, shape in enumerate(slide.shapes):
            if shape.shape_type == 13: # 13 is the constant for PICTURE
                # DEBUG: Print the raw XML attributes to see what is actually there
                try:
                    raw_descr = shape._element.nvPicPr.cNvPr.attrib.get("descr")
                    print(f"DEBUG: Slide {slide_num}, Shape {shape_idx}, descr attribute: '{raw_descr}'")
                except AttributeError:
                    print(f"DEBUG: Slide {slide_num}, Shape {shape_idx}, No descr attribute found.")

                # If raw_descr is None or empty, it's missing
                alt_text = getattr(shape, 'alt_text', None) # Fallback if library updated

                ##Check for None, empty string, or whitespace-only strings
                if raw_descr is None or not str(raw_descr).strip():
                     findings.append(Finding(
                        severity="serious",
                        rule_id="IMG_ALT",
                        wcag_criterion="1.1.1 Non-text Content",
                        message=f"Image on slide {slide_num} missing alt text",
                        slide_number=slide_num,
                        location=f"Slide {slide_num} - Shape {shape_idx}",
                        fix_hint="Right-click image → Format Picture → Alt Text."
                    ))
        return findings

    def _is_image(self, shape) -> bool:
        try:
            if getattr(shape, 'has_picture', False) or getattr(shape, 'image', None) is not None:
                return True
            if hasattr(shape, 'shape_type') and shape.shape_type in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.PICTURE_FRAME):
                return True
        except:
            pass
        return False

    def _check_slide_titles(self, slide, slide_num: int) -> List[Finding]:
        title = slide.shapes.title
        if not title or not title.text.strip():
            return [Finding(
                severity="critical",
                rule_id="SLIDE_TITLE",
                wcag_criterion="2.4.6 Headings and Labels",
                message=f"Slide {slide_num} has no title",
                slide_number=slide_num,
                location=f"Slide {slide_num}",
                fix_hint="Add a title placeholder or text box as the first element."
            )]
        return []

    def _check_empty_textboxes(self, slide, slide_num: int) -> List[Finding]:
        findings = []
        for idx, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                text = shape.text_frame.text.strip() if shape.text_frame.text else ""
                if text == "":
                    findings.append(Finding(
                        severity="minor",
                        rule_id="EMPTY_TEXTBOX",
                        wcag_criterion="1.3.1 Info and Relationships",
                        message=f"Empty text box found on slide {slide_num}",
                        slide_number=slide_num,
                        location=f"Slide {slide_num} - Shape {idx}",
                        fix_hint="Remove unused empty text boxes or add content."
                    ))
        return findings

    # ========================= NEW RULES =========================

    def _check_duplicate_titles(self, prs) -> List[Finding]:
        findings = []
        titles_seen: Dict[str, List[int]] = {}
        for slide_num, slide in enumerate(prs.slides, 1):
            title = slide.shapes.title
            if title and title.text.strip():
                clean = title.text.strip().lower()
                titles_seen.setdefault(clean, []).append(slide_num)

        for title, slides in titles_seen.items():
            if len(slides) > 1:
                findings.append(Finding(
                    severity="serious",
                    rule_id="DUPLICATE_TITLE",
                    wcag_criterion="2.4.6 Headings and Labels",
                    message=f"Duplicate slide title '{title}' on slides {slides}",
                    slide_number=slides[0],
                    location=f"Slides {slides}",
                    fix_hint="Make each slide title unique for better navigation."
                ))
        return findings

    def _check_reading_order(self, slide, slide_num: int) -> List[Finding]:
        if len(slide.shapes) < 3:
            return []
        return [Finding(
            severity="moderate",
            rule_id="READING_ORDER",
            wcag_criterion="1.3.2 Meaningful Sequence",
            message=f"Slide {slide_num}: Check reading order",
            slide_number=slide_num,
            location=f"Slide {slide_num}",
            fix_hint="Home → Arrange → Selection Pane. Drag items top to bottom in logical order.",
            confidence=0.6
        )]

    def _check_color_contrast(self, slide, slide_num: int) -> List[Finding]:
        findings = []
        for idx, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.color and hasattr(run.font.color, 'rgb') and run.font.color.rgb:
                        rgb = run.font.color.rgb
                        brightness = (rgb[0] * 299 + rgb[1] * 587 + rgb[2] * 114) / 1000
                        if brightness > 200:  # Very light text
                            findings.append(Finding(
                                severity="serious",
                                rule_id="LOW_CONTRAST",
                                wcag_criterion="1.4.3 Contrast (Minimum)",
                                message=f"Very light text on slide {slide_num} (possible low contrast)",
                                slide_number=slide_num,
                                location=f"Shape {idx}",
                                fix_hint="Use darker text. Aim for at least 4.5:1 contrast ratio."
                            ))
        return findings

    def _check_very_small_text(self, slide, slide_num: int) -> List[Finding]:
        findings = []
        for idx, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size and run.font.size.pt < 18:
                        findings.append(Finding(
                            severity="moderate",
                            rule_id="SMALL_TEXT",
                            wcag_criterion="1.4.4 Resize Text",
                            message=f"Very small text ({run.font.size.pt:.0f}pt) on slide {slide_num}",
                            slide_number=slide_num,
                            location=f"Shape {idx}",
                            fix_hint="Use font size ≥ 18pt for body text when possible."
                        ))
                        break
        return findings

    def _check_table_headers(self, slide, slide_num: int) -> List[Finding]:
        findings = []
        for idx, shape in enumerate(slide.shapes):
            if getattr(shape, 'has_table', False):
                table = shape.table
                if table.rows and len(table.rows) > 0:
                    if all(cell.text.strip() == '' for cell in table.rows[0].cells):
                        findings.append(Finding(
                            severity="serious",
                            rule_id="TABLE_HEADERS",
                            wcag_criterion="1.3.1 Info and Relationships",
                            message=f"Table on slide {slide_num} missing header row",
                            slide_number=slide_num,
                            location=f"Shape {idx}",
                            fix_hint="Select table → Table Design → Check 'Header Row'."
                        ))
        return findings

    def _check_merged_table_cells(self, slide, slide_num: int) -> List[Finding]:
        findings = []
        for idx, shape in enumerate(slide.shapes):
            if getattr(shape, 'has_table', False):
                table = shape.table
                for r_idx, row in enumerate(table.rows):
                    for c_idx, cell in enumerate(row.cells):
                        if getattr(cell, 'is_spanned', False):
                            findings.append(Finding(
                                severity="serious",
                                rule_id="TABLE_MERGED_CELLS",
                                wcag_criterion="1.3.1 Info and Relationships",
                                message=f"Merged cells in table on slide {slide_num}",
                                slide_number=slide_num,
                                location=f"Table {idx}, Row {r_idx+1}, Col {c_idx+1}",
                                fix_hint="Avoid merged cells in tables for screen reader compatibility."
                            ))
                            return findings  # Only flag once per table
        return findings

    def _check_image_only_slide(self, slide, slide_num: int) -> List[Finding]:
        has_text = any(s.has_text_frame and s.text_frame.text.strip() for s in slide.shapes if s != slide.shapes.title)
        has_image = any(self._is_image(s) or getattr(s, 'has_chart', False) for s in slide.shapes)

        if has_image and not has_text:
            return [Finding(
                severity="serious",
                rule_id="IMAGE_ONLY_SLIDE",
                wcag_criterion="1.1.1 Non-text Content",
                message=f"Slide {slide_num} is image-only with no supporting text",
                slide_number=slide_num,
                location=f"Slide {slide_num}",
                fix_hint="Add descriptive text explaining the visual content."
            )]
        return []

    def _check_chart_description(self, slide, slide_num: int) -> List[Finding]:
        findings = []
        for idx, shape in enumerate(slide.shapes):
            if getattr(shape, 'has_chart', False):
                alt_text = getattr(shape, 'alt_text', '') or ""
                if not str(alt_text).strip():
                    findings.append(Finding(
                        severity="serious",
                        rule_id="CHART_ALT",
                        wcag_criterion="1.1.1 Non-text Content",
                        message=f"Chart on slide {slide_num} missing description",
                        slide_number=slide_num,
                        location=f"Shape {idx}",
                        fix_hint="Add alt text describing the chart data and insights."
                    ))
        return findings

    def _check_text_embedded_in_image(self, slide, slide_num: int) -> List[Finding]:
        findings = []
        for idx, shape in enumerate(slide.shapes):
            if self._is_image(shape):
                name = str(getattr(shape, 'name', '')).lower()
                alt = str(getattr(shape, 'alt_text', '')).lower()
                keywords = ["screenshot", "diagram", "infographic", "text", "table", "chart"]
                if any(k in name or k in alt for k in keywords):
                    findings.append(Finding(
                        severity="moderate",
                        rule_id="EMBEDDED_TEXT_IMG",
                        wcag_criterion="1.4.5 Images of Text",
                        message=f"Image on slide {slide_num} likely contains text",
                        slide_number=slide_num,
                        location=f"Shape {idx}",
                        fix_hint="Extract text from the image into real text or provide full description in alt text."
                    ))
        return findings

    def _check_heading_structure(self, slide, slide_num: int) -> List[Finding]:
        # Simple heuristic for now
        if len([s for s in slide.shapes if s.has_text_frame]) > 5:
            return [Finding(
                severity="moderate",
                rule_id="HEADING_STRUCTURE",
                wcag_criterion="2.4.6 Headings and Labels",
                message=f"Slide {slide_num} has complex structure - verify heading hierarchy",
                slide_number=slide_num,
                location=f"Slide {slide_num}",
                fix_hint="Use consistent heading levels and logical structure."
            )]
        return []

    def _check_text_overlapping(self, slide, slide_num: int) -> List[Finding]:
        shapes = [s for s in slide.shapes if s.left is not None and s.top is not None]
        for i in range(len(shapes)):
            for j in range(i+1, len(shapes)):
                a, b = shapes[i], shapes[j]
                if (a.left < b.left + b.width and a.left + a.width > b.left and
                    a.top < b.top + b.height and a.top + a.height > b.top):
                    if a.has_text_frame or b.has_text_frame:
                        return [Finding(
                            severity="moderate",
                            rule_id="TEXT_OVERLAP",
                            wcag_criterion="1.3.2 Meaningful Sequence",
                            message=f"Overlapping content detected on slide {slide_num}",
                            slide_number=slide_num,
                            location=f"Slide {slide_num}",
                            fix_hint="Adjust positions so text doesn't overlap other objects."
                        )]
        return []